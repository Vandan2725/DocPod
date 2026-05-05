import warnings; warnings.filterwarnings("ignore")
from fastapi import APIRouter, Depends, HTTPException, UploadFile, File
from sqlalchemy.orm import Session
from database import get_db
from auth_utils import get_current_user
import models, os, uuid, threading

router = APIRouter()

ALLOWED_TYPES = {
    "application/pdf": "pdf",
    "application/msword": "doc",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/vnd.ms-powerpoint": "ppt",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
}
MAX_SIZE = 50 * 1024 * 1024


def extract_text(file_path: str, file_type: str) -> str:
    text = ""
    try:
        if file_type == "pdf":
            fitz = None
            try:
                import pymupdf as fitz
            except ImportError:
                try: import fitz
                except ImportError: pass
            if fitz:
                doc = fitz.open(file_path)
                parts = [p.get_text().strip() for p in doc if p.get_text().strip()]
                doc.close()
                text = "\n\n".join(parts)
        elif file_type in ("doc", "docx"):
            from docx import Document
            doc = Document(file_path)
            text = "\n".join(p.text.strip() for p in doc.paragraphs if p.text.strip())
        elif file_type in ("ppt", "pptx"):
            from pptx import Presentation
            prs = Presentation(file_path)
            parts = []
            for i, slide in enumerate(prs.slides):
                txts = [s.text.strip() for s in slide.shapes if hasattr(s, "text") and s.text.strip()]
                if txts: parts.append(f"Slide {i+1}. " + " ".join(txts))
            text = "\n\n".join(parts)
    except Exception as e:
        text = f"[Extraction error: {e}]"
    return text.strip()


def process_bg(doc_id: int, file_path: str, file_type: str, db_url: str):
    from sqlalchemy import create_engine
    from sqlalchemy.orm import sessionmaker
    eng = create_engine(db_url, connect_args={"check_same_thread": False} if "sqlite" in db_url else {})
    db = sessionmaker(bind=eng)()
    try:
        doc = db.query(models.Document).filter(models.Document.id == doc_id).first()
        if not doc: return
        doc.status = "extracting"; db.commit()
        text = extract_text(file_path, file_type)
        doc.extracted_text = text
        doc.status = "extracted" if text and not text.startswith("[") else "failed"
        db.commit()
    except Exception as e:
        d = db.query(models.Document).filter(models.Document.id == doc_id).first()
        if d: d.status = "failed"; d.extracted_text = str(e); db.commit()
    finally:
        db.close()


@router.post("/upload")
async def upload(
    file: UploadFile = File(...),
    db: Session = Depends(get_db),
    user: models.User = Depends(get_current_user)
):
    ext = file.filename.rsplit(".", 1)[-1].lower() if "." in file.filename else ""
    file_type = ALLOWED_TYPES.get(file.content_type or "") or {"pdf":"pdf","doc":"doc","docx":"docx","ppt":"ppt","pptx":"pptx"}.get(ext)
    if not file_type:
        raise HTTPException(400, "Unsupported type. Use PDF, DOC, DOCX, PPT, PPTX")
    contents = await file.read()
    if len(contents) > MAX_SIZE: raise HTTPException(400, "File too large. Max 50MB.")
    if len(contents) < 10:      raise HTTPException(400, "File is empty.")

    safe_name = f"{uuid.uuid4().hex}.{file_type}"
    path = os.path.join("uploads", safe_name)
    with open(path, "wb") as f: f.write(contents)

    doc = models.Document(user_id=user.id, filename=safe_name,
        original_filename=file.filename, file_type=file_type, file_size=len(contents))
    db.add(doc); db.commit(); db.refresh(doc)

    from database import DATABASE_URL
    threading.Thread(target=process_bg, args=(doc.id, path, file_type, DATABASE_URL), daemon=True).start()
    return {"document_id": doc.id, "filename": file.filename, "file_type": file_type, "status": "uploaded"}


@router.get("/{doc_id}")
def get_doc(doc_id: int, db: Session = Depends(get_db), user: models.User = Depends(get_current_user)):
    doc = db.query(models.Document).filter(models.Document.id == doc_id, models.Document.user_id == user.id).first()
    if not doc: raise HTTPException(404, "Not found")
    t = doc.extracted_text or ""
    return {"id": doc.id, "filename": doc.original_filename, "file_type": doc.file_type,
            "status": doc.status, "char_count": len(t),
            "extracted_text_preview": t[:300] + ("..." if len(t) > 300 else ""),
            "uploaded_at": doc.uploaded_at.isoformat()}


@router.get("/")
def list_docs(db: Session = Depends(get_db), user: models.User = Depends(get_current_user)):
    docs = db.query(models.Document).filter(models.Document.user_id == user.id).all()
    return [{"id": d.id, "filename": d.original_filename, "file_type": d.file_type,
             "status": d.status, "char_count": len(d.extracted_text or ""),
             "uploaded_at": d.uploaded_at.isoformat()} for d in docs]
